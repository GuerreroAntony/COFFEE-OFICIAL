from __future__ import annotations

import asyncio
import io
import logging
import re
import sys
from datetime import datetime, timedelta, timezone
from typing import Optional
from uuid import UUID

import httpx
from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, UploadFile, status

from app.config import settings
from app.database import execute_query, fetch_all, fetch_one
from app.dependencies import get_current_user
from app.schemas.materiais import (
    MaterialResponse,
    SyncStatusResponse,
    ToggleAIResponse,
)
from app.schemas.base import error_response, success_response
from app.services.embedding_service import generate_material_embeddings, remove_embeddings
from app.services.canvas_token_service import (
    CanvasAuthError,
    fetch_canvas_course_files,
    download_canvas_file,
)

logger = logging.getLogger("materiais")

router = APIRouter(prefix="/api/v1/materiais", tags=["materiais"])
disc_router = APIRouter(prefix="/api/v1/disciplinas", tags=["materiais"])


# ── Helpers ──────────────────────────────────────────────────

def _format_size(size_bytes: Optional[int]) -> Optional[str]:
    if not size_bytes:
        return None
    if size_bytes >= 1_000_000:
        return f"{size_bytes / 1_000_000:.1f} MB"
    return f"{size_bytes / 1_000:.0f} KB"


def _material_response(r) -> dict:
    return MaterialResponse(
        id=r["id"],
        disciplina_id=r["disciplina_id"],
        tipo=r["tipo"],
        nome=r["nome"],
        url_storage=r.get("url_storage"),
        fonte=r["fonte"],
        ai_enabled=r["ai_enabled"],
        size_bytes=r.get("size_bytes"),
        size_label=_format_size(r.get("size_bytes")),
        created_at=r["created_at"],
    ).model_dump(mode="json")


def _detect_tipo(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    mapping = {"pdf": "pdf", "pptx": "slide", "ppt": "slide", "docx": "documento", "doc": "documento", "jpg": "foto", "jpeg": "foto", "png": "foto"}
    return mapping.get(ext, "outro")


def _extract_content_type(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    mapping = {
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ppt": "application/vnd.ms-powerpoint",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "doc": "application/msword",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "png": "image/png",
    }
    return mapping.get(ext, "application/octet-stream")


async def _extract_text(content: bytes, filename: str) -> str:
    """Extract text from PDF or PPTX files."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    text = ""

    if ext == "pdf":
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            text = ""
    elif ext in ("pptx",):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            text = "\n".join(parts)
        except Exception:
            text = ""
    elif ext in ("docx", "doc"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        except Exception:
            text = ""

    return text.strip()


# ── GET /disciplina/{id} ─────────────────────────────────────

@disc_router.get("/{disciplina_id}/materiais")
async def listar_materiais(
    disciplina_id: UUID,
    user_id: UUID = Depends(get_current_user),
):
    """Lista materiais de uma disciplina."""
    enrolled = await fetch_one(
        "SELECT 1 FROM user_disciplinas WHERE user_id = $1 AND disciplina_id = $2",
        user_id, disciplina_id,
    )
    if not enrolled:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=error_response("ACCESS_DENIED", "Acesso negado"))

    rows = await fetch_all(
        """SELECT id, disciplina_id, tipo, nome, url_storage, fonte, ai_enabled, size_bytes, created_at
           FROM materiais WHERE disciplina_id = $1
           ORDER BY created_at DESC""",
        disciplina_id,
    )
    items = [_material_response(r) for r in rows]
    return success_response(items)


# ── GET /{id} ────────────────────────────────────────────────

@router.get("/{material_id}")
async def get_material(
    material_id: UUID,
    user_id: UUID = Depends(get_current_user),
):
    """Detalhe de um material."""
    row = await fetch_one(
        """SELECT m.id, m.disciplina_id, m.tipo, m.nome, m.url_storage,
                  m.fonte, m.ai_enabled, m.size_bytes, m.created_at
           FROM materiais m
           JOIN user_disciplinas ud ON ud.disciplina_id = m.disciplina_id AND ud.user_id = $2
           WHERE m.id = $1""",
        material_id, user_id,
    )
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=error_response("NOT_FOUND", "Material não encontrado"))
    return success_response(_material_response(row))


# ── PATCH /{id}/toggle-ai ───────────────────────────────────

@router.patch("/{material_id}/toggle-ai")
async def toggle_ai(
    material_id: UUID,
    background_tasks: BackgroundTasks,
    user_id: UUID = Depends(get_current_user),
):
    """Inverte ai_enabled e gera/remove embeddings."""
    row = await fetch_one(
        """UPDATE materiais m
           SET ai_enabled = NOT m.ai_enabled
           FROM user_disciplinas ud
           WHERE m.id = $1
             AND ud.disciplina_id = m.disciplina_id
             AND ud.user_id = $2
           RETURNING m.id, m.ai_enabled, m.texto_extraido, m.disciplina_id""",
        material_id, user_id,
    )
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=error_response("NOT_FOUND", "Material não encontrado"))

    # Generate or remove embeddings based on new state
    if row["ai_enabled"]:
        # ai_enabled became True → generate embeddings
        if row["texto_extraido"]:
            background_tasks.add_task(
                generate_material_embeddings,
                row["texto_extraido"], row["id"], row["disciplina_id"],
            )
    else:
        # ai_enabled became False → remove embeddings
        background_tasks.add_task(remove_embeddings, row["id"])

    resp = ToggleAIResponse(id=row["id"], ai_enabled=row["ai_enabled"])
    return success_response(resp.model_dump(mode="json"))


# ── PATCH /disciplina/{id}/materiais/enable-all-ai ────────────

@disc_router.patch("/{disciplina_id}/materiais/enable-all-ai")
async def enable_all_ai(
    disciplina_id: UUID,
    background_tasks: BackgroundTasks,
    user_id: UUID = Depends(get_current_user),
):
    """Enable AI for all materials of a discipline."""
    # Check enrollment
    enrolled = await fetch_one(
        "SELECT 1 FROM user_disciplinas WHERE user_id = $1 AND disciplina_id = $2",
        user_id, disciplina_id,
    )
    if not enrolled:
        raise HTTPException(status_code=403, detail=error_response("ACCESS_DENIED", "Acesso negado"))

    # Get materials that need AI enabled
    materials = await fetch_all(
        "SELECT id, texto_extraido FROM materiais WHERE disciplina_id = $1 AND ai_enabled = false",
        disciplina_id,
    )

    if not materials:
        return success_response({"updated_count": 0})

    # Enable AI for all
    await execute_query(
        "UPDATE materiais SET ai_enabled = true WHERE disciplina_id = $1 AND ai_enabled = false",
        disciplina_id,
    )

    # Generate embeddings for each material that has text
    for mat in materials:
        if mat["texto_extraido"]:
            background_tasks.add_task(
                generate_material_embeddings,
                mat["texto_extraido"],
                mat["id"],
                disciplina_id,
            )

    return success_response({"updated_count": len(materials)})


# ── POST /disciplina/{id}/materiais (manual upload) ──────────

@disc_router.post("/{disciplina_id}/materiais", status_code=status.HTTP_201_CREATED)
async def upload_material(
    disciplina_id: UUID,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    ai_enabled: bool = Form(True),
    user_id: UUID = Depends(get_current_user),
):
    """Upload manual de material."""
    enrolled = await fetch_one(
        "SELECT 1 FROM user_disciplinas WHERE user_id = $1 AND disciplina_id = $2",
        user_id, disciplina_id,
    )
    if not enrolled:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=error_response("ACCESS_DENIED", "Acesso negado"))

    content = await file.read()
    filename = file.filename or "arquivo"
    size_bytes = len(content)
    tipo = _detect_tipo(filename)
    content_type = _extract_content_type(filename)

    # Extract text
    texto_extraido = await _extract_text(content, filename)

    # Upload to Supabase Storage
    storage_path = f"{disciplina_id}/{filename}"
    upload_url = f"{settings.SUPABASE_URL}/storage/v1/object/materiais/{storage_path}"

    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.post(
            upload_url,
            content=content,
            headers={
                "Authorization": f"Bearer {settings.SUPABASE_KEY}",
                "Content-Type": content_type,
                "x-upsert": "true",
            },
        )
    if resp.status_code not in (200, 201):
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=error_response("AI_ERROR", "Erro no upload do arquivo"),
        )

    public_url = f"{settings.SUPABASE_URL}/storage/v1/object/public/materiais/{storage_path}"

    # Save to DB
    row = await fetch_one(
        """INSERT INTO materiais (disciplina_id, tipo, nome, url_storage, texto_extraido, fonte, ai_enabled, size_bytes)
           VALUES ($1, $2, $3, $4, $5, 'manual', $6, $7)
           RETURNING id, disciplina_id, tipo, nome, url_storage, fonte, ai_enabled, size_bytes, created_at""",
        disciplina_id, tipo, filename, public_url, texto_extraido, ai_enabled, size_bytes,
    )

    # Generate embeddings if ai_enabled and text extracted
    if ai_enabled and texto_extraido:
        background_tasks.add_task(
            generate_material_embeddings,
            texto_extraido, row["id"], disciplina_id,
        )

    return success_response(_material_response(row))


# ── POST /disciplina/{id}/sync ───────────────────────────────

@disc_router.post("/{disciplina_id}/sync")
async def trigger_sync(
    disciplina_id: UUID,
    background_tasks: BackgroundTasks,
    user_id: UUID = Depends(get_current_user),
):
    """Dispara sync de materiais do Canvas em background para uma disciplina."""
    enrolled = await fetch_one(
        "SELECT 1 FROM user_disciplinas WHERE user_id = $1 AND disciplina_id = $2",
        user_id, disciplina_id,
    )
    if not enrolled:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=error_response("ACCESS_DENIED", "Acesso negado"))

    disc = await fetch_one(
        "SELECT last_scraped_at, canvas_course_id FROM disciplinas WHERE id = $1",
        disciplina_id,
    )
    if not disc or not disc.get("canvas_course_id"):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=error_response("NO_CANVAS", "Disciplina sem canvas_course_id. Reconecte o ESPM."),
        )

    last_scraped = disc["last_scraped_at"]

    if last_scraped and (datetime.now(timezone.utc) - last_scraped) < timedelta(hours=settings.SYNC_COOLDOWN_HOURS):
        next_sync = last_scraped + timedelta(hours=settings.SYNC_COOLDOWN_HOURS)
        raise HTTPException(
            status_code=429,
            detail=error_response(
                "SYNC_COOLDOWN",
                "Sincronização em cooldown",
                extra={"next_sync_available_at": next_sync.isoformat()},
            ),
        )

    background_tasks.add_task(_sync_canvas_materials, disciplina_id, user_id)
    resp = SyncStatusResponse(status="triggered", last_synced_at=last_scraped)
    return success_response(resp.model_dump(mode="json"))


# ── POST /materiais/sync-all ─────────────────────────────────

@router.post("/sync-all")
async def trigger_sync_all(
    background_tasks: BackgroundTasks,
    user_id: UUID = Depends(get_current_user),
):
    """Dispara sync de materiais do Canvas para TODAS as disciplinas do usuário."""
    rows = await fetch_all(
        """SELECT d.id
           FROM disciplinas d
           JOIN user_disciplinas ud ON ud.disciplina_id = d.id
           WHERE ud.user_id = $1 AND d.canvas_course_id IS NOT NULL""",
        user_id,
    )
    if not rows:
        return success_response({"triggered": 0, "message": "Nenhuma disciplina com Canvas vinculado."})

    for row in rows:
        background_tasks.add_task(_sync_canvas_materials, row["id"], user_id)

    return success_response({"triggered": len(rows)})


# ── Canvas Materials Sync (background task) ──────────────────

# Pattern to auto-enable AI for class materials (e.g. "Aula 01", "Aula 12 - Tema")
_AULA_PATTERN = re.compile(r"(?i)\baula\s*\d+")

# Only sync these file types from Canvas (we can only extract text from PDF and PPTX)
_ALLOWED_SYNC_EXTENSIONS = {"pdf", "pptx", "ppt", "docx", "doc"}


async def _sync_canvas_materials(disciplina_id: UUID, user_id: UUID) -> None:
    """
    Fetch files from Canvas course, download new ones, store in Supabase,
    extract text, save to DB, and generate embeddings.
    """
    try:
        # 1. Get user's canvas_token
        user = await fetch_one(
            "SELECT canvas_token FROM users WHERE id = $1",
            user_id,
        )
        if not user or not user.get("canvas_token"):
            logger.error("[sync] user %s has no canvas_token", user_id)
            return

        canvas_token = user["canvas_token"]

        # 2. Get discipline's canvas_course_id
        disc = await fetch_one(
            "SELECT canvas_course_id FROM disciplinas WHERE id = $1",
            disciplina_id,
        )
        if not disc or not disc.get("canvas_course_id"):
            logger.error("[sync] disciplina %s has no canvas_course_id", disciplina_id)
            return

        canvas_course_id = disc["canvas_course_id"]

        # 3. Fetch file list from Canvas
        canvas_files = await fetch_canvas_course_files(canvas_token, canvas_course_id)
        logger.info("[sync] canvas course %d: %d files found", canvas_course_id, len(canvas_files))

        if not canvas_files:
            # Update last_scraped_at even if empty
            await execute_query(
                "UPDATE disciplinas SET last_scraped_at = NOW() WHERE id = $1",
                disciplina_id,
            )
            return

        # 4. Get existing canvas_file_ids for deduplication
        existing = await fetch_all(
            "SELECT canvas_file_id FROM materiais WHERE disciplina_id = $1 AND canvas_file_id IS NOT NULL",
            disciplina_id,
        )
        existing_ids = {r["canvas_file_id"] for r in existing}

        new_count = 0
        for cf in canvas_files:
            canvas_file_id = cf.get("id")
            if not canvas_file_id or canvas_file_id in existing_ids:
                continue

            filename = cf.get("display_name") or cf.get("filename") or "arquivo"
            file_url = cf.get("url")
            size_bytes = cf.get("size", 0)

            # Only sync PDF and PPTX files (skip images, docx, xlsx, zip, etc.)
            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
            if ext not in _ALLOWED_SYNC_EXTENSIONS:
                logger.debug("[sync] skipping %s (extension .%s not allowed)", filename, ext)
                continue

            if not file_url:
                logger.warning("[sync] file %s has no URL, skipping", filename)
                continue

            try:
                # 5. Download file from Canvas
                content = await download_canvas_file(canvas_token, file_url)
                logger.info("[sync] downloaded %s (%d bytes)", filename, len(content))

                # 6. Detect type and extract text
                tipo = _detect_tipo(filename)
                content_type = _extract_content_type(filename)
                texto_extraido = await _extract_text(content, filename)

                # 7. Upload to Supabase Storage
                storage_path = f"{disciplina_id}/{filename}"
                upload_url = f"{settings.SUPABASE_URL}/storage/v1/object/{settings.SUPABASE_STORAGE_BUCKET}/{storage_path}"

                async with httpx.AsyncClient(timeout=60.0) as client:
                    resp = await client.post(
                        upload_url,
                        content=content,
                        headers={
                            "Authorization": f"Bearer {settings.SUPABASE_KEY}",
                            "Content-Type": content_type,
                            "x-upsert": "true",
                        },
                    )
                if resp.status_code not in (200, 201):
                    logger.error("[sync] upload failed for %s: %s", filename, resp.status_code)
                    continue

                public_url = f"{settings.SUPABASE_URL}/storage/v1/object/public/{settings.SUPABASE_STORAGE_BUCKET}/{storage_path}"

                # 8. Auto-detect ai_enabled: "Aula XX" pattern → True
                ai_enabled = bool(_AULA_PATTERN.search(filename))

                # 9. Insert into DB
                row = await fetch_one(
                    """INSERT INTO materiais
                       (disciplina_id, tipo, nome, url_storage, texto_extraido, fonte, ai_enabled, size_bytes, canvas_file_id)
                       VALUES ($1, $2, $3, $4, $5, 'canvas', $6, $7, $8)
                       ON CONFLICT (canvas_file_id) WHERE canvas_file_id IS NOT NULL DO NOTHING
                       RETURNING id, disciplina_id, ai_enabled""",
                    disciplina_id, tipo, filename, public_url, texto_extraido,
                    ai_enabled, size_bytes, canvas_file_id,
                )

                if row:
                    new_count += 1
                    # 10. Generate embeddings if ai_enabled and text was extracted
                    if row["ai_enabled"] and texto_extraido:
                        try:
                            await generate_material_embeddings(
                                texto_extraido, row["id"], disciplina_id,
                            )
                        except Exception as emb_err:
                            logger.error("[sync] embedding error for %s: %s", filename, emb_err)

            except CanvasAuthError:
                logger.error("[sync] canvas token expired during sync for disciplina %s", disciplina_id)
                return  # Stop sync — token is dead
            except Exception as file_err:
                logger.error("[sync] error processing file %s: %s", filename, file_err)
                continue

        # 11. Update last_scraped_at
        await execute_query(
            "UPDATE disciplinas SET last_scraped_at = NOW() WHERE id = $1",
            disciplina_id,
        )

        logger.info("[sync] disciplina %s: %d new materials synced", disciplina_id, new_count)

    except Exception as e:
        logger.error("[sync] fatal error for disciplina %s: %s", disciplina_id, e)


async def sync_all_user_materials(user_id: UUID) -> None:
    """
    Sync materials for ALL user's disciplines with canvas_course_id.
    Called from ESPM connect after courses are upserted.
    """
    rows = await fetch_all(
        """SELECT d.id
           FROM disciplinas d
           JOIN user_disciplinas ud ON ud.disciplina_id = d.id
           WHERE ud.user_id = $1 AND d.canvas_course_id IS NOT NULL""",
        user_id,
    )
    if not rows:
        return

    logger.info("[sync-all] starting sync for %d disciplines (user %s)", len(rows), user_id)
    for row in rows:
        await _sync_canvas_materials(row["id"], user_id)
    logger.info("[sync-all] completed for user %s", user_id)
