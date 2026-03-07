from __future__ import annotations

import hashlib
import logging
import pathlib

logger = logging.getLogger(__name__)


def extract_text(filepath: str, file_type: str | None) -> str:
    """Extract text content from a file. Returns empty string on failure."""
    path = pathlib.Path(filepath)
    ext = (file_type or path.suffix.lstrip(".")).lower()

    try:
        if ext == "pdf":
            return _extract_pdf(path)
        elif ext in ("pptx", "ppt"):
            return _extract_pptx(path)
        else:
            logger.info("Unsupported file type '%s' for text extraction: %s", ext, path.name)
            return ""
    except Exception as e:
        logger.error("Text extraction failed for %s: %s", path.name, e)
        return ""


def _extract_pdf(path: pathlib.Path) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_pptx(path: pathlib.Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides.append("\n".join(texts))
    return "\n\n---\n\n".join(slides)


def compute_file_hash(filepath: str) -> str:
    """SHA-256 hash of file contents."""
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()
